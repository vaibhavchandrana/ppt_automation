from pathlib import Path
import contextlib
from pptx import Presentation   # importing the pptx library
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt   # importing point for setting font size
from pptx.dml.color import RGBColor  # importing RGB colors to set colors
# importing functions to create new charts
from pptx.chart.data import CategorySeriesData, CategoryChartData
# importing classes that raise custom exception
from Errors import ValidDigitError as ValidDigitHandler
from Errors import FeasebilityError as FeasebilityErrorHandler


# creating object of class ValidDigitHandler to handle the errors
obj_for_check_digit=ValidDigitHandler()


# creating object of class FeasebilityErrorHandler to handle the errors
obj_for_check_chart_data=FeasebilityErrorHandler()
        
# below function is used to validate the incoming data

def validateSecondSlideData(second_slide_data):
    """
    This function is used to check that incoming data is valid or not.
    Parameters:
    dictionery of data which has to be validated
    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True  
    try:
        # adding a validation variable to check data is valid or not
        
        for key, value in second_slide_data.items():
            # calling the checkType function to check the type of error
            validation = obj_for_check_digit.checkType(key, value)

    except Exception as e:
        validation=False
        print("Error is", e)

    return validation

def validateThirdSlideData(third_slide_data):
    """
    This function is used to check that incoming data is valid or not.

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True     # adding a validation variable to check data is valid or not
    try:
    # checking that is data recieved is valid or not
        for key, value in third_slide_data.items():
            # checking the key is not equal to Gender_Split_inside_store because Gender split ratio may contain string
            if key != "Gender_Split_inside_store" :
                # calling the checkType function of class ValidDigitHanler to check the  error
                if obj_for_check_digit.checkType(key, value) == False:
                    validation = False
    
    except Exception as e:
        validation=False
        print("Error is", e)
    return validation

def validateAndCalculateFourthSlideData(fourth_chart_data,passerby_vs_footfall):
    """
    This function is used to check that incoming data is valid or not.
    and it will also calculate the data for line chart

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    # validating the incoming data
    chart_data_validation = True
    try:
        for i in range(0, len(fourth_chart_data['passerby'])):
            # calling the CheckFeasibility function to check that data is feasibile or not
            chart_data_validation = obj_for_check_chart_data.CheckFeasibility(fourth_chart_data['passerby'][i],fourth_chart_data['footfall'][i])
            if chart_data_validation == True:
                passerby_vs_footfall[i] = round((fourth_chart_data['footfall'][i]/fourth_chart_data['passerby'][i])*100)
            else:
                break
    except Exception as e:
        chart_data_validation=False
        print("Error is", e)

    return chart_data_validation



# below function is used to update the slides
def updateSecondSlide(slides,second_slide_data):
    """
    This function is used to Update the data present in second slide.

    Parameters:
    1-List of all slides 
    2- Second slide data

    Returns:
    Nothing 
    Just save the changes in second slides
    """
# access the second slide where modification has to be done (0-indexed)
    second_slide = ppt.slides[1]

    # loop for traversing all the shapes
    for shape in second_slide.shapes:
        # checking that shape is equal to desired shaped

        if shape.name == "Rectangle 4":
            # fetching all the paragraph inside this component
            para = shape.text_frame.paragraphs

            # changing the paragraph with updated data
            para[3].text = f"Bring price point communication & Women window/easel stand – Improved passerby contribution by {second_slide_data['passerby_contribution']}% ( {second_slide_data['inc_in_customer']}% more customers stepped inside the store vs other stores)"
            # changing the font size of current paragraph
            para[3].font.size = Pt(12)

            # changing the paragraph with updated data
            para[4].text = f"Similarly, women walk-in’s contribution improved by {second_slide_data['women_walk_in_contribution']}% vs other stores"
            # changing the font size of current paragraph
            para[4].font.size = Pt(12)

            # changing the paragraph with updated data
            para[10].text = f"{second_slide_data['cust_notice_digital_window']}X customers noticed digital window vs static window"
            # changing the font size of current paragraph
            para[10].font.size = Pt(12)

            # changing the paragraph with updated data
            para[11].text = f"{second_slide_data['cust_walked_inside_digital']}% customers walked inside store after noticing digital window Vs {second_slide_data['cust_walked_inside_static']}% walked in after noticing static window"
            # changing the font size of current paragraph
            para[11].font.size = Pt(12)

            # changing the paragraph with updated data
            para[14].text = f"Only {second_slide_data['cust_notices_light_window']}% customers noticed light window, whereas {second_slide_data['cust_notices_static_window']}% customers notice static window"
            # changing the font size of current paragraph
            para[14].font.size = Pt(12)

    # saving the modified slide
    try:
        ppt.save('output.pptx')

    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")



def updateThirdSlide(slides,third_slide_data):
    """
    This function is used to Update the data present in third slide.

    Parameters:
    1-List of all slides 
    2- third slide data

    Returns:
    Nothing 
    Just save the changes in third slides
    """

    third_slide = ppt.slides[2]       # Getting the third slide to modify
    # loop for traversing all the shapes

    for shape in third_slide.shapes:
        # Condition for getting the desired shape( Component )

        if shape.name == "TextBox 8":
            # getting all the paragraph inside the box
            para = shape.text_frame.paragraphs

            # modifying the first paragraph
            para[0].text = f"{third_slide_data['passerby_vs_store_ff']}%" # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 15":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['Gender_Split_inside_store']}"  # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 22":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['Average_time_spent']} Min"  # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True    # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 36":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['Dwell_10_min']}%"   # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 51":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['Footfal_abandon']}%" # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True      # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 66":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['Dwell_time_vs_Bill_conversion']}%"  # Modifying the data
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True    # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "Text Placeholder 3":
            # changing the title of slide
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            para[0].text = f"{third_slide_data['market_footfall_walk_ins']}% market footfall walk-ins inside the stores, Men’s walk-ins is higher Average {third_slide_data['percentage_of_customer_spent']}% customers spent > 10 min inside the stores (potential customers)"  # Modifying the data of title
            # Changing the color of text size
            para[0].font.color.rgb = RGBColor(255, 0, 0)
            para[0].font.size = Pt(22)  # Changing the font size
            para[0].font.bold = True    # making the font bold

    # saving the third slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")


def updateForthSlide(slides,fourth_chart_data,passerby_vs_footfall):
    """
    This function is used to Update the charts present in forth slide.

    Parameters:
    1-List of all slides 
    2- forth slide data  (passerby and footfall)
    3- passerby vs footfall data

    Returns:
    Nothing 
    Just save the changes in forth slides
    """
    forth_slide = ppt.slides[3]        # accesing the fourth slide
    shape = forth_slide.shapes         # fetching all the shapes in the slide

    # getting the bar chart component from the shapes
    bar_chart = forth_slide.shapes[0].chart
    # getting the line chart component from the shapes
    line_chart = forth_slide.shapes[1].chart

    chart_data = ChartData()      # creating new instance of chart data

    # creating a instance for new chart
    new_bar_chart_data = CategoryChartData()

    # copying categoring from old charts
    new_bar_chart_data.categories = bar_chart.plots[0].categories

    # adding series of column chart with new data
    new_bar_series = new_bar_chart_data.add_series('Passer By', fourth_chart_data['passerby'])

    # adding series of column chart with new data
    new_bar_series = new_bar_chart_data.add_series('Footfall', fourth_chart_data['footfall'])

    # replacing the previous chart with new chart
    bar_chart.replace_data(new_bar_chart_data)

    # updating the line charts
    new_line_chart = CategoryChartData()

    # copying categoring from old charts
    new_line_chart.categories = line_chart.plots[0].categories

    # adding series of line chart with new data
    new_line_chart_series = new_line_chart.add_series('Vs Passer By', passerby_vs_footfall)

    # replacing the previous line chart with new chart
    line_chart.replace_data(new_line_chart)

    # saving the ppt into output pptx

    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")





# Open the existing PowerPoint presentation
ppt_file_path = Path('input.pptx')
try:
    if ppt_file_path.exists():
        ppt = Presentation(ppt_file_path)
        # Access all the slides in the presentation
        slides = ppt.slides


    # ============================ code for second slide start ================================

    # Data which has to be reflected in second slide
        second_slide_data = {
            "passerby_contribution": 15,
            "inc_in_customer": 20,
            "women_walk_in_contribution": 8,
            "cust_notice_digital_window": 3.5,
            "cust_walked_inside_digital": 14,
            "cust_walked_inside_static": 90,
            "cust_notices_light_window": 40,
            "cust_notices_static_window": 3
        }

        # checking that is data recieved is valid or not
        
        # calling validateSecondSlideData to validate the incoming data if it is true then only perform the below task
        second_slide_validation=validateSecondSlideData(second_slide_data)

        if second_slide_validation:
            ''' if data of second slide is valid then call the function updateSecondSlide to modify the second
            slide'''
            updateSecondSlide(slides, second_slide_data)
        else:
            print("Data of second slide is not valid")


    # ============================ code for third slide start ================================

        # Data which has to be reflected
        third_slide_data = {
            "passerby_vs_store_ff": 17.5,
            "Gender_Split_inside_store": "57:53",
            "Average_time_spent": 12.7,
            "Footfal_abandon": 18.3,
            "Dwell_10_min": 44.3,
            "Dwell_time_vs_Bill_conversion": 30.8,
            "market_footfall_walk_ins": 10,
            "percentage_of_customer_spent": 54
        }

        # calling validateThirdSlideData to validate the third slide data 
        third_slide_validation=validateThirdSlideData(third_slide_data)
        # perform below operation if no error is found
        if third_slide_validation:
            ''' if data of third slide is valid then call the function updateThirdSlide to modify the third
            slide'''
            updateThirdSlide(slides, third_slide_data)
        else:
            print("Data of third slide is not valid")


    # ============================ code for forth slide start ================================
        # data which has to modified in chart
        fourth_chart_data = {
            "passerby": [410689, 366210, 372079, 371521, 373346],
            "footfall": [73920, 55713, 52936, 50434, 48136]
        }
        
        # create a new list to store series data for drawing line chart
        passerby_vs_footfall = [0]*len(fourth_chart_data['passerby'])

        
        # if data is validated then perform below operation
        chart_data_validation=validateAndCalculateFourthSlideData(fourth_chart_data, passerby_vs_footfall)
        if chart_data_validation:
            ''' if data of forth slide data is valid then call the function updateForthSlide to modify the forth
         slide'''
            updateForthSlide(slides, fourth_chart_data, passerby_vs_footfall)
        else:
            print("Chart Data recieved is not valid ")


        #close the instance of pptx file    
        del ppt
    else:
        print('file path does not exist')
  
except Exception as e:
    print(e)
    print("Unable to open source file")
