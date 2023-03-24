from pptx import Presentation   # importing the pptx library
from pptx.chart.data import ChartData 
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt   # importing point for setting font size
from pptx.dml.color import RGBColor  # importing RGB colors to set colors
from pptx.chart.data import CategorySeriesData, CategoryChartData  # importing functions to create new charts

# Open the existing PowerPoint presentation
ppt = Presentation('input.pptx')

# Access the slides in the presentation
slides = ppt.slides

# access the second slide where modification has to be done (0-indexed)
slide = ppt.slides[1]

# Data which has to be reflected
first_data = {
    "passerby_contribution": 15,
    "inc_in_customer": 20,
    "women_walk_in_contribution": 8,
    "cust_notice_digital_window": 3.5,
    "cust_walked_inside_digital": 14,
    "cust_walked_inside_static": 90,
    "cust_notices_light_window": 40,
    "cust_notices_static_window": 3
}

# loop for traversing all the shapes
for shape in slide.shapes:
   # checking that shape is equal to desired shaped
    if shape.name == "Rectangle 4":
        para = shape.text_frame.paragraphs # fetching all the paragraph inside this component
        
        #changing the paragraph with updated data
        para[3].text = "Bring price point communication & Women window/easel stand – Improved passerby contribution by {}% ( {}% more customers stepped inside the store vs other stores)".format(
        first_data['passerby_contribution'], first_data['inc_in_customer'])
        para[3].font.size = Pt(12)    # changing the font size of current paragraph

        #changing the paragraph with updated data
        para[4].text = "Similarly, women walk-in’s contribution improved by {}% vs other stores".format(
                    first_data['women_walk_in_contribution'])
        para[4].font.size = Pt(12)    # changing the font size of current paragraph

        
        #changing the paragraph with updated data
        para[10].text = "{}X customers noticed digital window vs static window".format(
                    first_data['cust_notice_digital_window'])
        para[10].font.size = Pt(12)    # changing the font size of current paragraph


        #changing the paragraph with updated data
        para[11].text = "{}% customers walked inside store after noticing digital window Vs {}% walked in after noticing static window".format(
                    first_data['cust_walked_inside_digital'], first_data['cust_walked_inside_static'])
        para[11].font.size = Pt(12)    # changing the font size of current paragraph


        #changing the paragraph with updated data
        para[14].text = "Only {}% customers noticed light window, whereas {}% customers notice static window".format(
                    first_data['cust_notices_light_window'], first_data['cust_notices_static_window'])
        para[14].font.size = Pt(12)    # changing the font size of current paragraph


ppt.save('output.pptx') # saving the modified slide


slide = ppt.slides[2]       # Getting the third slide to modify

#Data which has to be reflected
data = {
    "passerby_vs_store_ff": 17.5,
    "Gender_Split_inside_store": "57:53",
    "Average_time_spent": 12.7,
    "Footfal_abandon": 18.3,
    "Dwell_10_min": 44.3,
    "Dwell_time_vs_Bill_conversion": 30.8,
    "market_footfall_walk_ins": 10,
    "customer_spent": 40
}


# loop for traversing all the shapes
for shape in slide.shapes:
    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 8":
        # getting all the paragraph inside the box
        para = shape.text_frame.paragraphs
        # modifying the first paragraph 
        para[0].text = "{}%".format(data['passerby_vs_store_ff'])  # Modifying the data
        para[0].font.size = Pt(25)   #Changing the font size
        para[0].font.bold = True   # making the font bold


    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 15":
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{}".format(data['Gender_Split_inside_store'])  # Modifying the data
        para[0].font.size = Pt(25)     #Changing the font size
        para[0].font.bold = True   # making the font bold


    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 22":
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{} Min".format(data['Average_time_spent']) # Modifying the data
        para[0].font.size = Pt(25) #Changing the font size
        para[0].font.bold = True    # making the font bold


    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 36":
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{}%".format(data['Dwell_10_min'])   # Modifying the data
        para[0].font.size = Pt(25)   #Changing the font size
        para[0].font.bold = True   # making the font bold

    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 51":
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{}%".format(data['Footfal_abandon'])  # Modifying the data
        para[0].font.size = Pt(25)     #Changing the font size
        para[0].font.bold = True      # making the font bold


    # Condition for getting the desired shape( Component )
    if shape.name == "TextBox 66":
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{}%".format(data['Dwell_time_vs_Bill_conversion']) # Modifying the data
        para[0].font.size = Pt(25)  #Changing the font size
        para[0].font.bold = True    # making the font bold


    # Condition for getting the desired shape( Component )
    if shape.name == "Text Placeholder 3":
        # changing the title of slide
        para = shape.text_frame.paragraphs
        # modifying the first paragraph
        para[0].text = "{}% market footfall walk-ins inside the stores, Men’s walk-ins is higher Average {}% customers spent >10 min inside the stores (potential customers)".format(
            data['market_footfall_walk_ins'], data['customer_spent'])   # Modifying the data of title 
        para[0].font.color.rgb = RGBColor(255, 0, 0)   #Changing the color of text size
        para[0].font.size = Pt(22) #Changing the font size
        para[0].font.bold = True    # making the font bold

ppt.save('output.pptx')    # saving the second slide 


slide = ppt.slides[3]        # accesing the fourth slide
shape = slide.shapes         # fetching all the shapes in the slide

chart = slide.shapes[0].chart    # getting the chart component from the shapes

chart_data = ChartData()      # creating new instance of chart data


# getting and storing detail of chart like series name and series values
for series in chart.series:
    chart_data.add_series(series.name, series.values)

# data which has to modified in chart
new_chart_data = {
    "new_values1": [100.0, 36.0, 37.0, 37.0, 37.0],
    "new_values2": [25.0, 155.0, 52.0, 50.0, 48.0],
    "new_values3": [4, 15.0, 14.0, 30.0, 12.0]   #pasesrby/footfall
}

# creating a instance for new chart
new_data = CategoryChartData()
new_data.categories = chart.plots[0].categories  # copying categoring from old charts
new_series = new_data.add_series('Passer By', new_chart_data['new_values1'])  # adding series of column chart with new data
new_series = new_data.add_series('Footfall', new_chart_data['new_values2'])  # adding series of column chart with new data
new_series = new_data.add_series('Vs Passer By', new_chart_data['new_values3'])  # adding series of line chart with new data


chart.replace_data(new_data)   # replacing the previous chart with new chart
ppt.save('output.pptx')     # saving the ppt into output pptx
