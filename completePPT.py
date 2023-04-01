from pathlib import Path
import contextlib
from pptx import Presentation   # importing the pptx library
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt   # importing point for setting font size
from pptx.dml.color import RGBColor  # importing RGB colors to set colors
# importing functions to create new charts
from pptx.chart.data import CategorySeriesData, CategoryChartData
# importing classes that raise custom exception
from Errors import ValidDigitError as ValidDigitHandler
from Errors import FeasebilityError as FeasebilityErrorHandler
import pptx


# creating object of class ValidDigitHandler to handle the errors
obj_for_check_digit = ValidDigitHandler()


# creating object of class FeasebilityErrorHandler to handle the errors
obj_for_check_chart_data = FeasebilityErrorHandler()

# below function is used to validate the incoming data


def validateSecondSlideData(second_slide_data):
    """
    This function is used to check that incoming data is valid or not.
    Parameters:
    dictionery of data which has to be validated
    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True
    try:
        # adding a validation variable to check data is valid or not

        for key, value in second_slide_data.items():
            # calling the checkType function to check the type of error
            validation = obj_for_check_digit.checkType(key, value)

    except Exception as e:
        validation = False
        print("Error is", e)

    return validation


def validateThirdSlideData(third_slide_data):
    """
    This function is used to check that incoming data is valid or not.
    Parameters:
    dictionery of data which has to be validated
    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True
    try:
        # adding a validation variable to check data is valid or not

        for key, value in third_slide_data.items():
            # calling the checkType function to check the type of error
            validation = obj_for_check_digit.checkType(key, value)

    except Exception as e:
        #if any exception happened then print the error
        validation = False
        print("Error is", e)

    return validation


def validateFourthSlideData(forth_slide_data):
    """
    This function is used to check that incoming data is valid or not.

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True     # adding a validation variable to check data is valid or not
    try:
        # checking that is data recieved is valid or not
        for key, value in forth_slide_data.items():
            # checking the key is not equal to Gender_Split_inside_store because Gender split ratio may contain string
            if key != "Gender_Split_inside_store":
                # calling the checkType function of class ValidDigitHanler to check the  error
                if obj_for_check_digit.checkType(key, value) == False:
                    validation = False

    except Exception as e:
        validation = False
        print("Error is", e)
    return validation

def validateElevenData(eleven_slide_data):
    """
    This function is used to check that incoming data is valid or not.
    Parameters:
    dictionery of data which has to be validated
    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True
    try:
        # adding a validation variable to check data is valid or not

        for key, value in eleven_slide_data.items():
            # calling the checkType function to check the type of error
            if key=='Passer_Vs_Store_Walk_ins_Gender_wise' or key=='instore_Walk_ins_Gender_wise':
                continue
            else:
                validation = obj_for_check_digit.checkType(key, value)

    except Exception as e:
        validation = False
        print("Error is", e)

    return validation

def validateData(seventh_slide_data):
    """
    This function is used to check that incoming data is valid or not.

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    validation = True     # adding a validation variable to check data is valid or not
    try:
        # checking that is data recieved is valid or not
        for key, value in seventh_slide_data.items():
            # calling the checkType function of class ValidDigitHanler to check the  error
            if obj_for_check_digit.checkType(key, value) == False:
                validation = False

    except Exception as e:
        validation = False
        print("Error is", e)
    return validation

def validateAndCalculateFifthSlideData(fourth_chart_data, passerby_vs_footfall):
    """
    This function is used to check that incoming data is valid or not.
    and it will also calculate the data for line chart

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    # validating the incoming data
    chart_data_validation = True
    try:
        for i in range(0, len(fourth_chart_data['passerby'])):
            # calling the CheckFeasibility function to check that data is feasibile or not
            chart_data_validation = obj_for_check_chart_data.CheckFeasibility(
                fourth_chart_data['passerby'][i], fourth_chart_data['footfall'][i])
            if chart_data_validation == True:
                passerby_vs_footfall[i] = round(
                    (fourth_chart_data['footfall'][i]/fourth_chart_data['passerby'][i])*100)
                passerby_vs_footfall[i]/=100
            else:
                break
    except Exception as e:
        chart_data_validation = False
        print("Error is", e)

    return chart_data_validation


def validateAndCalculateSixthSlideData(sixth_chart_data,sixth_line_data ):
    """
    This function is used to check that incoming data is valid or not.
    and it will also calculate the data for line chart

    Parameters:
    dictionery of data which has to be validated

    Returns:
    true : if data is validated
    false: if data is invalid
    """
    # validating the incoming data
    chart_data_validation = True
    try:
        for i in range(0, len(sixth_chart_data['Customer_spent'])):
            # calling the CheckFeasibility function to check that data is feasibile or not
            chart_data_validation = obj_for_check_chart_data.CheckFeasibilitySixth(
                sixth_chart_data['Customer_spent'][i], sixth_chart_data['No_of_bills'][i])
            if chart_data_validation == True:
                sixth_line_data[i] = round(
                    (sixth_chart_data['No_of_bills'][i]/sixth_chart_data['Customer_spent'][i])*100)
                sixth_line_data[i]/=100
            else:
                break
    except Exception as e:
        chart_data_validation = False
        print("Error is", e)

    return chart_data_validation

def calculatePercentage(data):
    '''this function is used to calculate percentage for charts data
    input: List of integers
    return: Nothing just replace data with its percentage
    '''
    for i in range(0,len(data)):
        data[i]/=100

# below function is used to update the slides
def updateSecondSlide(slides, second_slide_data,month_of_data):
    """
    This function is used to Update the data present in second slide.

    Parameters:
    1-List of all slides 
    2- Second slide data

    Returns:
    Nothing 
    Just save the changes in second slides
    """
# access the second slide where modification has to be done (0-indexed)
    second_slide = ppt.slides[1]

    # loop for traversing all the shapes
    for shape in second_slide.shapes:
        # checking that shape is equal to desired shaped

        # updating the month in second slide
        if shape.name == "Star: 6 Points 20":
            # fetching all the paragraph inside this component
            para = shape.text_frame.paragraphs[0]
            para.text=month_of_data
            para.font.size=Pt(12)
            para.font.bold=True

        if shape.name == "Star: 6 Points 18":
            # fetching all the paragraph inside this component
            para = shape.text_frame.paragraphs[0]
            para.text=month_of_data
            para.font.size=Pt(12)
            para.font.bold=True

        if shape.name == "Rectangle 4":
            # fetching all the paragraph inside this component
            para = shape.text_frame.paragraphs

            # changing the paragraph with updated data
            para[
                3].text = f"Bring price point communication & Women window/easel stand – Improved passerby contribution by {second_slide_data['passerby_contribution']}% ( {second_slide_data['inc_in_customer']}% more customers stepped inside the store vs other stores)"
            # changing the font size of current paragraph
            para[3].font.size = Pt(12)

            # changing the paragraph with updated data
            para[4].text = f"Similarly, women walk-in’s contribution improved by {second_slide_data['women_walk_in_contribution']}% vs other stores"
            # changing the font size of current paragraph
            para[4].font.size = Pt(12)

            # changing the paragraph with updated data
            para[10].text = f"{second_slide_data['cust_notice_digital_window']}X customers noticed digital window vs static window"
            # changing the font size of current paragraph
            para[10].font.size = Pt(12)

            # changing the paragraph with updated data
            para[11].text = f"{second_slide_data['cust_walked_inside_digital']}% customers walked inside store after noticing digital window Vs {second_slide_data['cust_walked_inside_static']}% walked in after noticing static window"
            # changing the font size of current paragraph
            para[11].font.size = Pt(12)

            # changing the paragraph with updated data
            para[14].text = f"Only {second_slide_data['cust_notices_light_window']}% customers noticed light window, whereas {second_slide_data['cust_notices_static_window']}% customers notice static window"
            # changing the font size of current paragraph
            para[14].font.size = Pt(12)

    # saving the modified slide
    try:
        ppt.save('output.pptx')

    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")

def updateThirdSlide(slides,third_slide_data,month_of_data):
    """
    This function is used to Update the data present in third slide.

    Parameters:
    1-List of all slides 
    2- Third slide data

    Returns:
    Nothing 
    Just save the changes in second slides
    """
    
    third_slide = ppt.slides[2]
    for shape in third_slide.shapes:
        # Updating the month of slide
        if shape.name == "Star: 6 Points 16":
            # fetching all the paragraph inside this component
            para = shape.text_frame.paragraphs[0]
            para.text=month_of_data
            para.font.size=Pt(12)
            para.font.bold=True

        if shape.name == "Rectangle 4":
            para = shape.text_frame.paragraphs
            para[2].text=f"{third_slide_data['customers_moving_Sneaker_wall_vs_total_footfall']}% customers moving towards Sneaker wall vs total footfall walk inside the stores"
            para[2].font.size = Pt(12)

            para[7].text=f"{third_slide_data['customers_greeted_by_staff']}% of customers greeted by staff, whereas staff engagement with customers is during purchase journey is {third_slide_data['staff_engagement_with_customers_during_purchase']}% vs overall walk-ins"
            para[7].font.size = Pt(12)

            para[8].text=f"Staff engagement during customer journey drops to {third_slide_data['Staff_engagement_during_customer_journey']}% on weekends"
            para[8].font.size = Pt(12)
            
            para[9].text=f"Hushpuppies Abmi mall store has high engagement % : {third_slide_data['Hushpuppies_Abmi_mall_store_engagement']}% vs overall walk-ins"
            para[9].font.size = Pt(12)

            para[10].text=f"Vegas mall Bata engagement % : {third_slide_data['Vegas_mall_Bata_engagement']}% vs overall walk-ins"
            para[10].font.size = Pt(12)

            para[11].text=f"Similarly, Ambi mall Bata engagement % : {third_slide_data['Ambi_mall_Bata_engagement']}% vs overall walk-ins"
            para[11].font.size = Pt(12)

            para[12].text=f"Hypothesis : Is more engagement led to more conversion ? Its stand true, as Vegas mall conversion is {third_slide_data['Vegas_mall_conversion']}% Vs {third_slide_data['Ambi_mall_conversion']}% conversion on Ambi mall."
            para[12].font.size = Pt(12)

            para[14].text=f" {third_slide_data['footwear_picked_up_customers']}X times footwear picked up customers"
            para[14].font.size = Pt(12)

            para[15].text=f"More Serious/interested shoppers walked in during weekends, same is reflecting no. of shoes picked from shelf {third_slide_data['no_of_shoes_picked_from_shelf']}X times vs weekdays"
            para[15].font.size = Pt(12)

            para[17].text=f" {third_slide_data['customers_reaches_till_1st_half']}% customers reaches till 1st half of the store"
            para[17].font.size = Pt(12)

            para[18].text=f"Only {third_slide_data['customers_look_at_side_panels']}% customers look at side panels"
            para[18].font.size = Pt(12)

    # saving the forth slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")


def updateForthSlide(slides, fourth_slide_data):
    """
    This function is used to Update the data present in third slide.

    Parameters:
    1-List of all slides 
    2- third slide data

    Returns:
    Nothing 
    Just save the changes in third slides
    """

    forth_slide = ppt.slides[3]       # Getting the third slide to modify
    # loop for traversing all the shapes

    for shape in forth_slide.shapes:
        # Condition for getting the desired shape( Component )

        if shape.name == "TextBox 8":
            # getting all the paragraph inside the box
            para = shape.text_frame.paragraphs

            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['passerby_vs_store_ff']}%"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 15":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['Gender_Split_inside_store']}"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 22":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['Average_time_spent']} Min"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True    # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 36":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['Dwell_10_min']}%"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True   # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 51":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['Footfal_abandon']}%"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True      # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "TextBox 66":
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data
            para[0].text = f"{fourth_slide_data['Dwell_time_vs_Bill_conversion']}%"
            para[0].font.size = Pt(25)  # Changing the font size
            para[0].font.bold = True    # making the font bold

        # Condition for getting the desired shape( Component )
        if shape.name == "Text Placeholder 3":
            # changing the title of slide
            para = shape.text_frame.paragraphs
            # modifying the first paragraph
            # Modifying the data of title
            para[
                0].text = f"{fourth_slide_data['market_footfall_walk_ins']}% market footfall walk-ins inside the stores, Men’s walk-ins is higher Average {fourth_slide_data['percentage_of_customer_spent']}% customers spent > 10 min inside the stores (potential customers)"
            # Changing the color of text size
            para[0].font.color.rgb = RGBColor(255, 0, 0)
            para[0].font.size = Pt(22)  # Changing the font size
            para[0].font.bold = True    # making the font bold

    # saving the forth slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")


def updateFifthSlide(slides, fifth_slide_first_chart_data,fifth_first_line_chart,fifth_slide_second_line_data,fifth_slide_third_chart_data,fifth_slide_third_line_chart,fifth_slide_forth_chart_data,fifth_slide_forth_line_chart):
    """
    This function is used to Update the data present in third slide.

    Parameters:
    1-List of all slides 
    2- fifth slide data for all four charts

    Returns:
    Nothing 
    Just save the changes in third slides
    """
    fifth_slide=ppt.slides[4] 

    for shape in fifth_slide.shapes:
        
        if shape.name == "Chart 4":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Passer By', fifth_slide_first_chart_data['passerby'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('Footfall', fifth_slide_first_chart_data['footfall'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart
        
        if shape.name=="Chart 5":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Vs Passer By', fifth_first_line_chart)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        ''' updating the second chart of slide'''
        if shape.name=="Chart 18":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Vs Passer By', fifth_slide_second_line_data)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        if shape.name == "Chart 6":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Passer By', fifth_slide_third_chart_data['passerby'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('Footfall', fifth_slide_third_chart_data['footfall'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart
        
        '''updatating the line chart of third chart '''
        if shape.name=="Chart 7":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Vs Passer By', fifth_slide_third_line_chart)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        # '''updating the forth chart bar graph '''
        if shape.name == "Chart 17":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Passer By', fifth_slide_forth_chart_data['passerby'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('Footfall', fifth_slide_forth_chart_data['footfall'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart
        
        # '''updatating the line chart of forth chart '''
        if shape.name=="Chart 8":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Vs Passer By', fifth_slide_forth_line_chart)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

    
    # saving the forth slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")

def updateSixthSlideData(slides,sixth_slide_first_chart_bar_data,sixth_slide_first_line_chart,sixth_slide_second_chart_bar_data,sixth_slide_second_line_chart,sixth_slide_third_chart_bar_data,sixth_slide_third_line_chart_data,sixth_slide_data):
    """
    This function is used to Update the data present in sixth slide.

    Parameters:
    1-List of all slides 
    2- fifth slide data for all four charts

    Returns:
    Nothing 
    Just save the changes in third slides
    """
    fifth_slide=ppt.slides[5] 

    for shape in fifth_slide.shapes:

        
        if shape.name == "Rectangle 4":
                para = shape.text_frame.paragraphs[0]
                para.text=f"{sixth_slide_data['converted_into_sale']}% People converted into Sale"
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold=True
                para.font.size=Pt(15)
        
        if shape.name == "Rectangle 5":
                para = shape.text_frame.paragraphs[0]
                para.text=f"{sixth_slide_data['Average_time_spent']} min"
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold=True
                para.font.size=Pt(15)

        if shape.name == "Rectangle 10":
                para = shape.text_frame.paragraphs[1]
                para.text=f"{sixth_slide_data['dwell_more_than_10']}% "
                para.font.color.rgb = RGBColor(255, 255, 255)
                para.font.bold=True
                para.font.size=Pt(15)
        # '''updatating the second chart '''
        if shape.name == "Chart 7":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Customer spent > 10 min', sixth_slide_second_chart_bar_data['Customer_spent'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('No of bills', sixth_slide_second_chart_bar_data['No_of_bills'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart
        
        if shape.name=="Chart 8":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Customer spent Vs No of bills', sixth_slide_second_line_chart)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        # ''' updating the first chart of slide'''
        if shape.name == "Chart 9":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Customer spent > 10 min', sixth_slide_first_chart_bar_data['Customer_spent'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('No of bills', sixth_slide_first_chart_bar_data['No_of_bills'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart
        

        ''' updating the line chart of first  chart of slide'''
        if shape.name=="Chart 11":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Customer spent Vs No of bills', sixth_slide_first_line_chart)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart
        

        # ''' updating the third chart of slide'''
        if shape.name == "Chart 13":
            bar_chart = shape.chart 
            new_first_data = CategoryChartData()
            new_first_data.categories = bar_chart.plots[0].categories  # copying categoring from old charts
            new_first_series = new_first_data.add_series('Customer spent > 10 min', sixth_slide_third_chart_bar_data['Customer_spent'])  # adding series of column chart with new data
            new_first_series = new_first_data.add_series('No of bills', sixth_slide_third_chart_bar_data['No_of_bills'])  # adding series of column chart with new data
            bar_chart.replace_data(new_first_data)   # replacing the previous chart with new chart

        ''' updating the third chart of first  chart of slide'''
        if shape.name=="Chart 15":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Customer spent Vs No of bills', sixth_slide_third_line_chart_data)  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart
    
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")


def updateSeventhSlide(slides,seventh_slide_data):
    """
        This function is used to Update the data present in sixth slide.

        Parameters:
        1-List of all slides 
        2- fifth slide data for all four charts

        Returns:
        Nothing 
        Just save the changes in third slides
        """
    seventh_slide=ppt.slides[6]  
    i=0
    for shape in seventh_slide.shapes:
        
    
        # # if shape.shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
        if shape.name=="Text Placeholder 2":
            paras=shape.text_frame.paragraphs
            paras[0].text=f"{seventh_slide_data['shoes_picked_up_from_Shelfs']}X times shoes picked up from Shelfs, staff engaged with only {seventh_slide_data['staff_engaged_with_customer']}% customers vs overall walk-ins"
            paras[0].font.size=Pt(20)

        if i==6:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Customers_Moved_towards_Shelf']}%"
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.bold=True
            para.font.size=Pt(25)

        if i==9:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['No_of_time_shoes_Picked_from_Shelf']}X"
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.bold=True
            para.font.size=Pt(25)

        if i==11:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Customers_greeted_by_staff']}%"
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.bold=True
            para.font.size=Pt(25)

        if i==17:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['No_of_time_shoes_Picked_from_Shelf_weekday']}X"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)
        
        
        if i==18:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['No_of_time_shoes_Picked_from_Shelf_weekend']}X"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)


        if i==19:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Customers_greeted_by_staff_ambi_bata']}%"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)

        
        
        if i==20:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Customers_greeted_by_staff_vegas']}%"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)

        
        if i==22:
            para=shape.text_frame.paragraphs[0]
            para.text=f"Hypothesis : Is more engagement led to more conversion ? Its stand true, as Vegas mall conversion is {seventh_slide_data['vega_mall_conversion']}% Vs {seventh_slide_data['ambi_mall_conversion']}% conversion on Ambi mall. "
            para.font.size=Pt(15)
        
        
        if i==26:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Staff_engagement_Vs_overall_walkins']}%"
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.bold=True
            para.font.size=Pt(25)


        if i==33:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Customers_greeted_by_staff_ambi_hp']}%"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)
        
        
        if i==37:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Staff_engagement_Vs_overall_walkins_ambi_bata']}%"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)


        if i==38:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Staff_engagement_Vs_overall_walkins_vegas']}%"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)


        if i==39:
            para=shape.text_frame.paragraphs[0]
            para.text=f"{seventh_slide_data['Staff_engagement_Vs_overall_ambi_hp']} %"
            para.font.color.rgb = RGBColor(12, 99, 27)
            para.font.bold=True
            para.font.size=Pt(14)

        i+=1
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")

def updateEightSlide(slides,eight_slide_data,eight_slide_chart_data,table_data_ambi_mall,table_data_vegas_mall):
    """
        This function is used to Update the data present in sixth slide.

        Parameters:
        1-List of all slides 
        2- fifth slide data for all four charts

        Returns:
        Nothing 
        Just save the changes in third slides
        """
    eight_slide=ppt.slides[7]     
    for shape in eight_slide.shapes:
        if shape.name == "Table 7":
            table = shape.table
            k = 0
            for row in table.rows:
                if k==2:
                    cell =row.cells
                    cell[1].text_frame.paragraphs[0].text=f"{table_data_ambi_mall['people_noticed_static']}%"
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[2].text_frame.paragraphs[0].text=f"{table_data_ambi_mall['noticed_vs_stepped_inside_static']}%"
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[3].text_frame.paragraphs[0].text=f"{table_data_ambi_mall['people_noticed_digital']}%"
                    cell[3].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[4].text_frame.paragraphs[0].text=f"{table_data_ambi_mall['noticed_vs_stepped_inside_digital']}%"
                    cell[4].text_frame.paragraphs[0].font.size=Pt(12)
                
                if k==3:
                    cell =row.cells
                    cell[1].text_frame.paragraphs[0].text=f"{table_data_vegas_mall['people_noticed_static']}%"
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[2].text_frame.paragraphs[0].text=f"{table_data_vegas_mall['noticed_vs_stepped_inside_static']}%"
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[3].text_frame.paragraphs[0].text=f"{table_data_vegas_mall['people_noticed_digital']}%"
                    cell[3].text_frame.paragraphs[0].font.size=Pt(12)
                    cell[4].text_frame.paragraphs[0].text=f"{table_data_vegas_mall['noticed_vs_stepped_inside_digital']}%"
                    cell[4].text_frame.paragraphs[0].font.size=Pt(12)
                
                k+=1

        if shape.name == 'Chart 2':
            bar_chart = shape.chart
            new_first_data = CategoryChartData()
            # copying categoring from old charts
            new_first_data.categories = bar_chart.plots[0].categories
            # adding series of column chart with new data
            new_first_series = new_first_data.add_series(
                'Static Window', eight_slide_chart_data['static_window'])
            # adding series of column chart with new data
            new_first_series = new_first_data.add_series(
                'Digital Window', eight_slide_chart_data['digital_window'])
            # adding series of column chart with new data
            new_first_series = new_first_data.add_series(
                'Light Window', eight_slide_chart_data['light_window'])
            # replacing the previous chart with new chart
            bar_chart.replace_data(new_first_data)

        if shape.name == "Text Placeholder 6":
            para = shape.text_frame.paragraphs[0]
            para.text = f"{eight_slide_data['customer_notice_digital_vs_static']}X customers noticed digital window vs static window, {eight_slide_data['higher_passerby_ratio']}X higher passer by ratio"
            para.font.size = Pt(20)

    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")

def updateNinthSlide(slides,ninth_slide_data,ninth_slide_first_chart,ninth_slide_second_chart):
    """
        This function is used to Update the data present in sixth slide.

        Parameters:
        1-List of all slides 
        2- fifth slide data for all four charts

        Returns:
        Nothing 
        Just save the changes in third slides
        """
    ninth_slide=ppt.slides[8] 
    
    change_in_second_data=[round(ninth_slide_second_chart['other_store'][1]*100-ninth_slide_second_chart['other_store'][0]*100,1),round(ninth_slide_second_chart['innovative_window_store'][1]*100-ninth_slide_second_chart['innovative_window_store'][0]*100,1)]


    change_in_first_data=[
    round(ninth_slide_first_chart['other_store'][1]*100 -ninth_slide_first_chart['other_store'][0]*100,1),round(ninth_slide_first_chart['innovative_window_store'][1]*100-ninth_slide_first_chart['innovative_window_store'][0]*100,1)]

    for shape in ninth_slide.shapes:

        if shape.name == "Text Placeholder 2":
            para=shape.text_frame.paragraphs[0]
            para.text=f"{ninth_slide_data['success_rate']}% success rate, out of 2 stores, one store noticed +{ninth_slide_data['increase_in_passerby_store']}% increase in passer by share"
            para.font.size=Pt(20)
            
            

        if shape.name == "Chart 5":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Other Store', ninth_slide_first_chart['other_store'])  # adding series of line chart with new data
            new_line_series = new_line_chart.add_series('Innovative Window Store', ninth_slide_first_chart['innovative_window_store'])  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        if shape.name == "Chart 7":
            line_chart=shape.chart
            new_line_chart=CategoryChartData()
            new_line_chart.categories = line_chart.plots[0].categories  # copying categoring from old charts
            new_line_series = new_line_chart.add_series('Other Store', ninth_slide_second_chart['other_store'])  # adding series of line chart with new data
            new_line_series = new_line_chart.add_series('Innovative Window Store', ninth_slide_second_chart['innovative_window_store'])  # adding series of line chart with new data
            line_chart.replace_data(new_line_chart)   # replacing the previous line chart with new chart

        if shape.name=="Rectangle 10":
            para = shape.text_frame.paragraphs[0]
            if change_in_first_data[0]>0:
                para.text=f"+{change_in_first_data[0]}%"
                
            else:
                para.text=f"{change_in_first_data[0]}%"
            para.font.color.rgb = RGBColor( 255,0, 0)
            para.font.bold=True
            para.font.size=Pt(15)

        if shape.name=="Rectangle 34":
            para = shape.text_frame.paragraphs[0]
            if change_in_first_data[1]>0:
                para.text=f"+{change_in_first_data[1]}%"
                
            else:
                para.text=f"{change_in_first_data[1]}%"
            para.font.color.rgb = RGBColor( 255,0, 0)
            para.font.bold=True
            para.font.size=Pt(15)
        
        if shape.name=="Rectangle 35":
            para = shape.text_frame.paragraphs[0]
            if change_in_second_data[0]>0:
                para.text=f"+{change_in_second_data[0]}%"
                
            else:
                para.text=f"{change_in_second_data[0]}%"
            para.font.color.rgb = RGBColor( 255,0, 0)
            para.font.bold=True
            para.font.size=Pt(15)

        if shape.name=="Rectangle 36":
            para = shape.text_frame.paragraphs[0]
            if change_in_second_data[1]>0:
                para.text=f"+{change_in_second_data[1]}%"
                
            else:
                para.text=f"{change_in_second_data[1]}%"
            para.font.color.rgb = RGBColor( 255,0, 0)
            para.font.bold=True
            para.font.size=Pt(15)

        
        if shape.name=="Rectangle 35":
            para = shape.text_frame.paragraphs[0]



def updateEleventhSlide(chika_compound,bandra_matrix):
    """
    This function is used to Update the data present in eleventh slide.

    Parameters:
    1-List of all slides 
    2- fifth slide data for all four charts

    Returns:
    Nothing 
    Just save the changes in third slides
    """
    
    eleven_slide=ppt.slides[10] 
    for shape in eleven_slide.shapes:
        if shape.name == "Table 1":
            table = shape.table
            r=0 # take a variable which keep track of row number
            for row in table.rows:
                if r == 1: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=f"{chika_compound['Passer_Vs_Store_Walkins']}%"    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=f"{bandra_matrix['Passer_Vs_Store_Walkins']}%"    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                if r == 2: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=chika_compound['Passer_Vs_Store_Walk_ins_Gender_wise']    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=bandra_matrix['Passer_Vs_Store_Walk_ins_Gender_wise']    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width


                if r == 3: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=chika_compound['instore_Walk_ins_Gender_wise']    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=bandra_matrix['instore_Walk_ins_Gender_wise']    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width


                if r == 4: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=f"{chika_compound['Avg_time_Spend_inside_Store']} Min"    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=f"{bandra_matrix['Avg_time_Spend_inside_Store']} min"    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width


                if r == 5: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=f"{chika_compound['Customers_spend_less_than_2_min']}%"    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=f"{bandra_matrix['Customers_spend_less_than_2_min']}%"    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width


                if r == 6: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=f"{chika_compound['Potential_customers_10_min']}%"    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=f"{bandra_matrix['Potential_customers_10_min']}%"    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width


                if r == 7: # frtching the specified row from table
                    cell =row.cells # feting all the columns from row
                    cell[1].text_frame.paragraphs[0].text=f"{chika_compound['Potential_customers_conversion']}%"    # updating text of column 1
                    cell[1].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                    cell[2].text_frame.paragraphs[0].text=f"{bandra_matrix['Potential_customers_conversion']}%"    # updating text of column 2
                    cell[2].text_frame.paragraphs[0].font.size=Pt(12) # seting font size to specified width

                r+=1
    # saving the forth slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")

def updateTenthSlide(slides,tenth_slide_data):
    """
    This function is used to Update the data present in elventh slide.

    Parameters:
    1-List of all slides 
    2- fifth slide data for all four charts

    Returns:
    Nothing 
    Just save the changes in third slides
    """
    i = 0
    tenth_slide=ppt.slides[9] 
    for shape in tenth_slide.shapes:
        if i == 0:
            img_file_path = "Heatmap.jpg"
            try:
                if img_file_path:
                    new_image_file = img_file_path
                    # create new image part from new image file
                    new_pptx_img = pptx.parts.image.Image.from_file(new_image_file)

                    #  to figure out what image you're actually changing...
                    img_shape = shape

                    # get part and rId from shape we need to change
                    slide_part, rId = img_shape.part, img_shape._element.blip_rId
                    image_part = slide_part.related_part(rId)
                    
                    # overwrite old blob info with new blob info
                    image_part.blob = new_pptx_img._blob
            except Exception as e:
                print("Image file not found ",e)
        if i==4:
            para=shape.text_frame.paragraphs
            para[2].text=f"{tenth_slide_data['customers_reaches_till_1st_half']}% customers reaches till 1st half of the store"
            para[2].font.size=Pt(14)
            para[2].font.color.rgb = RGBColor(255, 0, 0)


            para[4].text=f"Only {tenth_slide_data['customers_walk_at_end_store']}% customers walk at end of the store"
            para[4].font.size=Pt(14)
            para[4].font.color.rgb = RGBColor(255, 0, 0)

            para[6].text=f"Only {tenth_slide_data['customers_looked_side_panels']}% customers look at side panels"
            para[6].font.size=Pt(14)
            para[6].font.color.rgb = RGBColor(255, 0, 0)

            para[8].text=f"Center tables / Gondola's catches the interest of {tenth_slide_data['customers_at_central_tables']}% customers"
            para[8].font.size=Pt(14)
            para[8].font.color.rgb = RGBColor(255, 0, 0)

        i+=1
    # saving the tenth slide
    try:
        ppt.save('output.pptx')  # saving the modified slide
    except Exception as e:
        print(e)
        print("File is already opened . Please close the file first to write")




# Open the existing PowerPoint presentation
ppt_file_path = Path('input.pptx')
try:
    if ppt_file_path.exists():
        ppt = Presentation(ppt_file_path)
        # Access all the slides in the presentation
        slides = ppt.slides

    # ============================ code for second slide start ================================
        month_of_data="Sept'22"
        # Data which has to be reflected in second slide
        second_slide_data = {
            "passerby_contribution": 15,
            "inc_in_customer": 20,
            "women_walk_in_contribution": 8,
            "cust_notice_digital_window": 3.5,
            "cust_walked_inside_digital": 14,
            "cust_walked_inside_static": 90,
            "cust_notices_light_window": 40,
            "cust_notices_static_window": 3
        }

        # checking that is data recieved is valid or not

        # calling validateSecondSlideData to validate the incoming data if it is true then only perform the below task
        second_slide_validation = validateSecondSlideData(second_slide_data)

        if second_slide_validation:
            ''' if data of second slide is valid then call the function updateSecondSlide to modify the second
            slide'''
            updateSecondSlide(slides, second_slide_data,month_of_data)
        else:
            print("Data of second slide is not valid")

    # ============================ code for third slide start ================================
    # Data which has to be reflected in third slide
        third_slide_data = {
            'customers_moving_Sneaker_wall_vs_total_footfall': 42,
            'customers_greeted_by_staff': 93,
            'staff_engagement_with_customers_during_purchase': 12,
            'Staff_engagement_during_customer_journey': 9,
            'Hushpuppies_Abmi_mall_store_engagement': 21,
            'Vegas_mall_Bata_engagement': 13.3,
            'Ambi_mall_Bata_engagement': 7.3,
            'Vegas_mall_conversion': 17,
            'Ambi_mall_conversion': 8,
            'footwear_picked_up_customers': 2.5,
            'no_of_shoes_picked_from_shelf': 1.7,
            'customers_reaches_till_1st_half': 52,
            'customers_look_at_side_panels': 33
        }
        # calling validateThirdSlideData to validate the third slide data
        third_slide_validation = validateThirdSlideData(third_slide_data)
        # perform below operation if no error is found
        if third_slide_validation:
            ''' if data of third slide is valid then call the function updateThirdSlide to modify the third
            slide'''
            updateThirdSlide(slides, third_slide_data,month_of_data)
        else:
            print("Data of third slide is not valid")

    # ============================ code for forth slide start ================================

        # Data which has to be reflected
        fourth_slide_data = {
            "passerby_vs_store_ff": 16.5,
            "Gender_Split_inside_store": "57:53",
            "Average_time_spent": 12.7,
            "Footfal_abandon": 18.3,
            "Dwell_10_min": 44.3,
            "Dwell_time_vs_Bill_conversion": 30.8,
            "market_footfall_walk_ins": 10,
            "percentage_of_customer_spent": 54
        }

        # calling validateThirdSlideData to validate the forth slide data
        forth_slide_validation = validateFourthSlideData(fourth_slide_data)
        # perform below operation if no error is found
        if forth_slide_validation:
            ''' if data of fourth slide is valid then call the function updateThirdSlide to modify the third
            slide'''
            updateForthSlide(slides, fourth_slide_data)
        else:
            print("Data of fourth slide is not valid")

    # ============================ code for fifth slide start================================

        # data which has to modified in first bar chart
        fifth_slide_first_chart_data = {
            "passerby": [410800, 366210, 372079, 371521, 373346],
            "footfall": [73920, 55713, 52936, 50434, 48136]
        }
        # data which has to modified in first  linechart
        fifth_first_line_chart=[0,0,0,0,0]

        # data which has to modified in second line chart
        fifth_slide_second_line_data=[12.1,12.3,12.2,12.4,12.1]

        # data which has to modified in third bar chart
        fifth_slide_third_chart_data = {
            "passerby": [295114.0, 534452.0, 575753.0, 288526.0],
            "footfall": [52195.0, 144539.0, 51004.0, 33401.0]
        }

        # data which has to modified in third line chart
        fifth_slide_third_line_chart=[0, 0, 0, 0]

        # data which has to modified in forth bar chart
        fifth_slide_forth_chart_data = {
            "passerby": [410689.0, 366000.0, 372000.0],
            "footfall": [73920.0, 55713.0, 52900.0]
        }
        fifth_slide_forth_line_chart=[0, 0, 0]
        
        # here validating the data for first chart of fifth slide
        fifth_validation_first=validateAndCalculateFifthSlideData(fifth_slide_first_chart_data, fifth_first_line_chart)

        # here validating the data for third chart of fifth slide
        fifth_validation_third=validateAndCalculateFifthSlideData(fifth_slide_third_chart_data,fifth_slide_third_line_chart)

        # here validating the data for forth chart of fifth slide
        fifth_validation_forth=validateAndCalculateFifthSlideData(fifth_slide_forth_chart_data,fifth_slide_forth_line_chart)
        
# checking that the data of fifth_slide is valid or not
        if fifth_validation_first and fifth_validation_third and fifth_validation_forth:
            #if it is valid then call updateFifthSlideData to update value of slide
            updateFifthSlide(slides, fifth_slide_first_chart_data,fifth_first_line_chart,fifth_slide_second_line_data,fifth_slide_third_chart_data,fifth_slide_third_line_chart,fifth_slide_forth_chart_data,fifth_slide_forth_line_chart)

        else:
            print("data of fifth slide is not valid")


    # ============================ code for sixth slide start================================
        # Data which has to be reflected in sixth slide

        sixth_slide_data = {
            'Average_time_spent': 13.5,
            'dwell_more_than_10': 43.1,
            'converted_into_sale': 30
        }

        sixth_slide_first_chart_bar_data = {
            'Customer_spent': [23900.0, 52700.0, 23215.0, 15445.0],
            'No_of_bills': [10167.0, 7900.0, 8464.0, 7451.0]
        }
        sixth_slide_first_line_chart = [0, 0, 0, 0]

        sixth_slide_second_line_chart = [0, 0, 0]
        sixth_slide_second_chart_bar_data = {
            'Customer_spent': [54434.0, 60200.0, 23375.0],
            'No_of_bills': [15002.0, 18893.0, 1859.0]
        }

        sixth_slide_third_chart_bar_data = {
            'Customer_spent': [72247.0, 42449.0],
            'No_of_bills': [22600.0, 11300.0]
        }
        sixth_slide_third_line_chart_data = [0, 0]

        # here validating the data for first chart of sixth slide

        sixth_slide_first_validation=validateAndCalculateSixthSlideData(sixth_slide_first_chart_bar_data,sixth_slide_first_line_chart)

                # here validating the data for second chart of sixth slide
        sixth_slide_second_validation=validateAndCalculateSixthSlideData(sixth_slide_second_chart_bar_data,sixth_slide_second_line_chart)

                # here validating the data for third chart of sixth slide
        sixth_slide_third_validation=validateAndCalculateSixthSlideData(sixth_slide_third_chart_bar_data,sixth_slide_third_line_chart_data)
       
        if sixth_slide_first_validation and sixth_slide_second_validation and sixth_slide_third_validation:
            ''' checking that the data is valid or not 
            if it is valid the call updateSixthSlideData to update the data in sixth slide'''
            
            updateSixthSlideData(slides,sixth_slide_first_chart_bar_data,sixth_slide_first_line_chart,sixth_slide_second_chart_bar_data,sixth_slide_second_line_chart,sixth_slide_third_chart_bar_data,sixth_slide_third_line_chart_data,sixth_slide_data)
        
        else:
            print("data of second slide is not valid")


        # ============================ code for seventh slide start================================
     
    # Data which has to be reflected in seventh slide

        seventh_slide_data={
        'shoes_picked_up_from_Shelfs':3.4,
        'staff_engaged_with_customer':11.6,
        'Customers_Moved_towards_Shelf':51,
        'No_of_time_shoes_Picked_from_Shelf':2.3,
        'No_of_time_shoes_Picked_from_Shelf_weekday':2.3,
        'No_of_time_shoes_Picked_from_Shelf_weekend':2.4,
        'Customers_greeted_by_staff':93,
        'Customers_greeted_by_staff_ambi_bata':83,
        'Customers_greeted_by_staff_vegas':85,
        'Customers_greeted_by_staff_ambi_hp':91,
        'vega_mall_conversion':17,
        'ambi_mall_conversion':9,
        'Staff_engagement_Vs_overall_walkins':11.6,
        'Staff_engagement_Vs_overall_walkins_ambi_bata':7.3,
        'Staff_engagement_Vs_overall_walkins_vegas':13.5,
        'Staff_engagement_Vs_overall_ambi_hp':20,
        }

        # here validating the data of seventh slide
        if validateData(seventh_slide_data):
            updateSeventhSlide(slides, seventh_slide_data)
        else:
            print("data of seventh slide is not valid")

 # ============================ code for eight slide start================================
    # Data which has to be reflected in eighth slide
        eight_slide_data = {
            'customer_notice_digital_vs_static': 5.3,
            'higher_passerby_ratio': 1.5
        }

        eight_slide_chart_data = {
            'static_window': [3, 9],
            'digital_window': [10, 14],
            'light_window': [4, 16]
        }
        # converting the chart data into percentage
        calculatePercentage(eight_slide_chart_data['static_window'])
        # converting the chart data into percentage
        calculatePercentage(eight_slide_chart_data['digital_window'])
        # converting the chart data into percentage
        calculatePercentage(eight_slide_chart_data['light_window'])

        # table data for ambi mall       
        table_data_ambi_mall={
            "people_noticed_static":2,
            "noticed_vs_stepped_inside_static":10,
            "people_noticed_digital":11,
            "noticed_vs_stepped_inside_digital":14,
        }

        # table data for ambi mall       
        table_data_vegas_mall={
            "people_noticed_static":3,
            "noticed_vs_stepped_inside_static":7,
            "people_noticed_digital":9,
            "noticed_vs_stepped_inside_digital":15,
        }
        # here validating the eighth slide data
        if validateData(eight_slide_data) and validateData(table_data_ambi_mall) and validateData(table_data_vegas_mall):
            # if data is valid then call updateEightSlide function 
            updateEightSlide(slides,eight_slide_data,eight_slide_chart_data,table_data_ambi_mall,table_data_vegas_mall)

 # ============================ code for ninth slide start================================
        # Data which has to be reflected in eighth slide
        ninth_slide_data={
            'success_rate':45,
            'increase_in_passerby_store':1.7
        }

        ninth_slide_first_chart={
            'other_store':[6.5, 6.9],
            'innovative_window_store':[3.6, 5.3]
        }
        # calculating percentage for ninth_slide_first_chart data
        calculatePercentage(ninth_slide_first_chart['other_store'])
        calculatePercentage(ninth_slide_first_chart['innovative_window_store'])

        ninth_slide_second_chart={
            'other_store':[6.8, 7.2],
            'innovative_window_store':[18.1, 18.3]
        }
        # calculating percentage for ninth_slide_second_chart data
        calculatePercentage(ninth_slide_second_chart['other_store'])
        calculatePercentage(ninth_slide_second_chart['innovative_window_store'])
        
        updateNinthSlide(slides,ninth_slide_data,ninth_slide_first_chart,ninth_slide_second_chart)
        

 # ============================ code for tenth slide start================================
        tenth_slide_data = {
            'customers_reaches_till_1st_half': 53,
            'customers_walk_at_end_store': 10,
            'customers_looked_side_panels': 33,
            'customers_at_central_tables': 66
        }
        updateTenthSlide(slides,tenth_slide_data)


# ============================ code for eleventh slide start================================

        chika_compound={
            'Passer_Vs_Store_Walkins':18,
            'Passer_Vs_Store_Walk_ins_Gender_wise':'20:15', 
            'instore_Walk_ins_Gender_wise':'59:41',
            'Avg_time_Spend_inside_Store':12.7,
            'Customers_spend_less_than_2_min':9,
            'Potential_customers_10_min':45,
            'Potential_customers_conversion':55
        }

        bandra_matrix={
            'Passer_Vs_Store_Walkins':8,
            'Passer_Vs_Store_Walk_ins_Gender_wise':'07:08', 
            'instore_Walk_ins_Gender_wise':'40:60',
            'Avg_time_Spend_inside_Store':13.2,
            'Customers_spend_less_than_2_min':7,
            'Potential_customers_10_min':50,
            'Potential_customers_conversion':50 
        }
         
        eleventh_data_validation_one=validateElevenData(chika_compound)
        eleventh_data_validation_two=validateElevenData(bandra_matrix)
        
        # Checking 11th data is validated or not
        if eleventh_data_validation_one and eleventh_data_validation_two:
            updateEleventhSlide(chika_compound,bandra_matrix)


        del ppt
        print("Presentation updated successfully")
    else:
        print('file path does not exist')


except Exception as e:
    print(e)
    print("Unable to open source file")
